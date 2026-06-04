"""PPTX slide title reader."""

from __future__ import annotations

from pathlib import Path
from typing import List


class PptxStructureReader:
    @staticmethod
    def read_slide_titles(path: Path) -> List[str]:
        try:
            from pptx import Presentation
        except Exception:
            return []

        presentation = Presentation(str(path))
        titles: List[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title and getattr(slide.shapes.title, "text", ""):
                title = slide.shapes.title.text.strip()
            titles.append(title or f"Slide {index}")
        return titles
